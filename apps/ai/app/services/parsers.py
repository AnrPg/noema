# =============================================================================
# DOCUMENT PARSERS
# =============================================================================

import io
from typing import Dict

import structlog

logger = structlog.get_logger()


class BaseParser:
    """Base class for all document parsers."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse document content."""
        raise NotImplementedError


class PDFParser(BaseParser):
    """Parser for PDF documents."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse PDF document."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed, returning empty result")
            return self._empty_result(filename, "pdf")

        doc = fitz.open(stream=content, filetype="pdf")

        sections = []
        full_text = []
        word_count = 0

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = str(page.get_text())
            word_count += len(text.split())
            sections.append(
                {
                    "title": f"Page {page_num + 1}",
                    "content": text,
                    "page": page_num + 1,
                    "section_type": "text",
                    "metadata": {},
                }
            )
            full_text.append(text)

        doc.close()

        return {
            "filename": filename,
            "file_type": "pdf",
            "sections": sections,
            "total_pages": len(sections),
            "word_count": word_count,
            "metadata": {},
        }

    def _empty_result(self, filename: str, file_type: str) -> Dict:
        return {
            "filename": filename,
            "file_type": file_type,
            "sections": [],
            "total_pages": 0,
            "word_count": 0,
            "metadata": {},
        }


class DocxParser(BaseParser):
    """Parser for DOCX documents."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse DOCX document."""
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx not installed, returning empty result")
            return self._empty_result(filename, "docx")

        doc = Document(io.BytesIO(content))

        sections = []
        word_count = 0
        current_section = {"title": None, "content": "", "section_type": "text", "metadata": {}}

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Check if it's a heading
            style_name = para.style.name if para.style else None
            if style_name and style_name.startswith("Heading"):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {
                    "title": text,
                    "content": "",
                    "section_type": "heading",
                    "metadata": {"style": style_name},
                }
            else:
                current_section["content"] += text + "\n"
                word_count += len(text.split())

        if current_section["content"]:
            sections.append(current_section)

        return {
            "filename": filename,
            "file_type": "docx",
            "sections": sections,
            "total_pages": None,
            "word_count": word_count,
            "metadata": {},
        }

    def _empty_result(self, filename: str, file_type: str) -> Dict:
        return {
            "filename": filename,
            "file_type": file_type,
            "sections": [],
            "total_pages": None,
            "word_count": 0,
            "metadata": {},
        }


class TextParser(BaseParser):
    """Parser for plain text files."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse plain text file."""
        text = content.decode("utf-8", errors="ignore")

        # Split by double newlines to create sections
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

        sections = []
        word_count = 0

        for para in paragraphs:
            sections.append(
                {
                    "title": None,
                    "content": para,
                    "page": None,
                    "section_type": "text",
                    "metadata": {},
                }
            )
            word_count += len(para.split())

        return {
            "filename": filename,
            "file_type": "txt",
            "sections": sections,
            "total_pages": None,
            "word_count": word_count,
            "metadata": {},
        }


class MarkdownParser(BaseParser):
    """Parser for Markdown files."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse Markdown file."""
        text = content.decode("utf-8", errors="ignore")

        sections = []
        word_count = 0
        current_section = {"title": None, "content": "", "section_type": "text", "metadata": {}}

        for line in text.split("\n"):
            # Check for headers
            if line.startswith("#"):
                if current_section["content"]:
                    sections.append(current_section)

                # Count header level
                level = 0
                for char in line:
                    if char == "#":
                        level += 1
                    else:
                        break

                title = line.lstrip("#").strip()
                current_section = {
                    "title": title,
                    "content": "",
                    "section_type": f"h{level}",
                    "metadata": {"level": level},
                }
            else:
                current_section["content"] += line + "\n"
                word_count += len(line.split())

        if current_section["content"]:
            sections.append(current_section)

        return {
            "filename": filename,
            "file_type": "md",
            "sections": sections,
            "total_pages": None,
            "word_count": word_count,
            "metadata": {},
        }


class HTMLParser(BaseParser):
    """Parser for HTML files."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse HTML file."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            # Fallback to basic text extraction
            text = content.decode("utf-8", errors="ignore")
            import re

            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()

            return {
                "filename": filename,
                "file_type": "html",
                "sections": [
                    {"title": None, "content": text, "section_type": "text", "metadata": {}}
                ],
                "total_pages": None,
                "word_count": len(text.split()),
                "metadata": {},
            }

        soup = BeautifulSoup(content, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style"]):
            element.decompose()

        sections = []
        word_count = 0

        # Extract headings and their content
        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p"]):
            text = tag.get_text(strip=True)
            if text:
                tag_name = getattr(tag, "name", None) or "p"
                is_heading = isinstance(tag_name, str) and tag_name.startswith("h")
                section_type = tag_name if is_heading else "text"
                sections.append(
                    {
                        "title": text if is_heading else None,
                        "content": text if not is_heading else "",
                        "section_type": section_type,
                        "metadata": {},
                    }
                )
                word_count += len(text.split())

        return {
            "filename": filename,
            "file_type": "html",
            "sections": sections,
            "total_pages": None,
            "word_count": word_count,
            "metadata": {},
        }


class ExcelParser(BaseParser):
    """Parser for Excel files."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse Excel file."""
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl not installed, returning empty result")
            return self._empty_result(filename, "xlsx")

        workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True)

        sections = []
        word_count = 0

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                rows.append(" | ".join(row_text))

            content_text = "\n".join(rows)
            word_count += len(content_text.split())

            sections.append(
                {
                    "title": sheet_name,
                    "content": content_text,
                    "section_type": "table",
                    "metadata": {"sheet": sheet_name},
                }
            )

        workbook.close()

        return {
            "filename": filename,
            "file_type": "xlsx",
            "sections": sections,
            "total_pages": None,
            "word_count": word_count,
            "metadata": {"sheets": workbook.sheetnames},
        }

    def _empty_result(self, filename: str, file_type: str) -> Dict:
        return {
            "filename": filename,
            "file_type": file_type,
            "sections": [],
            "total_pages": None,
            "word_count": 0,
            "metadata": {},
        }


class PowerPointParser(BaseParser):
    """Parser for PowerPoint files."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        extract_images: bool = False,
        extract_tables: bool = True,
    ) -> Dict:
        """Parse PowerPoint file."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed, returning empty result")
            return self._empty_result(filename, "pptx")

        prs = Presentation(io.BytesIO(content))

        sections = []
        word_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    text = text.strip()
                    if text:
                        slide_text.append(text)
                        word_count += len(text.split())

            sections.append(
                {
                    "title": f"Slide {slide_num}",
                    "content": "\n".join(slide_text),
                    "page": slide_num,
                    "section_type": "slide",
                    "metadata": {},
                }
            )

        return {
            "filename": filename,
            "file_type": "pptx",
            "sections": sections,
            "total_pages": len(sections),
            "word_count": word_count,
            "metadata": {},
        }

    def _empty_result(self, filename: str, file_type: str) -> Dict:
        return {
            "filename": filename,
            "file_type": file_type,
            "sections": [],
            "total_pages": 0,
            "word_count": 0,
            "metadata": {},
        }


# Export all parsers
__all__ = [
    "BaseParser",
    "PDFParser",
    "DocxParser",
    "TextParser",
    "MarkdownParser",
    "HTMLParser",
    "ExcelParser",
    "PowerPointParser",
]
