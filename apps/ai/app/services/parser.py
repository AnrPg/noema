# =============================================================================
# DOCUMENT PARSER SERVICE
# =============================================================================

import io
import re
from pathlib import Path
from typing import Dict, List

import structlog

logger = structlog.get_logger()


class DocumentParser:
    """
    Service for parsing various document formats and extracting text content.
    """

    SUPPORTED_FORMATS = {
        "pdf": ["application/pdf"],
        "docx": ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"],
        "pptx": ["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
        "txt": ["text/plain"],
        "md": ["text/markdown", "text/x-markdown"],
        "html": ["text/html"],
        "image": ["image/png", "image/jpeg", "image/gif", "image/webp"],
    }

    def __init__(self):
        self._ocr_enabled = False
        self._check_dependencies()

    def _check_dependencies(self):
        """Check which parsing libraries are available."""
        try:
            import pytesseract

            self._ocr_enabled = True
        except ImportError:
            logger.warning("pytesseract not available, OCR disabled")

    async def parse(
        self,
        content: bytes,
        filename: str,
        content_type: str,
    ) -> Dict:
        """
        Parse document and extract text content.

        Args:
            content: File content as bytes
            filename: Original filename
            content_type: MIME type

        Returns:
            Dict with extracted text, metadata, and structure
        """
        ext = Path(filename).suffix.lower().lstrip(".")

        parsers = {
            "pdf": self._parse_pdf,
            "docx": self._parse_docx,
            "pptx": self._parse_pptx,
            "txt": self._parse_text,
            "md": self._parse_markdown,
            "html": self._parse_html,
            "png": self._parse_image,
            "jpg": self._parse_image,
            "jpeg": self._parse_image,
            "gif": self._parse_image,
            "webp": self._parse_image,
        }

        parser = parsers.get(ext)
        if not parser:
            raise ValueError(f"Unsupported file format: {ext}")

        return await parser(content, filename)

    async def _parse_pdf(self, content: bytes, filename: str) -> Dict:
        """Parse PDF document."""
        import fitz  # PyMuPDF

        doc = fitz.open(stream=content, filetype="pdf")

        pages = []
        full_text = []
        images = []

        for page_num, page in enumerate(doc):
            text = page.get_text()
            pages.append(
                {
                    "page": page_num + 1,
                    "text": text,
                    "word_count": len(text.split()),
                }
            )
            full_text.append(text)

            # Extract images for OCR if needed
            for img_index, img in enumerate(page.get_images()):
                xref = img[0]
                base_image = doc.extract_image(xref)
                images.append(
                    {
                        "page": page_num + 1,
                        "index": img_index,
                        "format": base_image["ext"],
                    }
                )

        return {
            "text": "\n\n".join(full_text),
            "pages": pages,
            "metadata": {
                "title": doc.metadata.get("title", filename),
                "author": doc.metadata.get("author", ""),
                "page_count": len(doc),
                "word_count": sum(p["word_count"] for p in pages),
                "has_images": len(images) > 0,
            },
            "structure": self._extract_structure(full_text),
        }

    async def _parse_docx(self, content: bytes, filename: str) -> Dict:
        """Parse Word document."""
        from docx import Document

        doc = Document(io.BytesIO(content))

        paragraphs = []
        full_text = []
        headings = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            paragraphs.append(
                {
                    "text": text,
                    "style": para.style.name if para.style else "Normal",
                }
            )
            full_text.append(text)

            if para.style and "Heading" in para.style.name:
                headings.append(
                    {
                        "text": text,
                        "level": self._get_heading_level(para.style.name),
                    }
                )

        # Extract tables
        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                rows.append(cells)
            tables.append(rows)

        return {
            "text": "\n\n".join(full_text),
            "paragraphs": paragraphs,
            "headings": headings,
            "tables": tables,
            "metadata": {
                "word_count": len(" ".join(full_text).split()),
                "paragraph_count": len(paragraphs),
                "table_count": len(tables),
            },
            "structure": self._build_structure_from_headings(headings),
        }

    async def _parse_pptx(self, content: bytes, filename: str) -> Dict:
        """Parse PowerPoint presentation."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        slides = []
        full_text = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            shapes_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    shapes_text.append(
                        {
                            "type": shape.shape_type.name
                            if hasattr(shape.shape_type, "name")
                            else str(shape.shape_type),
                            "text": shape.text,
                        }
                    )

            combined_text = "\n".join(slide_text)
            slides.append(
                {
                    "slide": slide_num + 1,
                    "text": combined_text,
                    "shapes": shapes_text,
                }
            )
            full_text.append(combined_text)

        return {
            "text": "\n\n".join(full_text),
            "slides": slides,
            "metadata": {
                "slide_count": len(slides),
                "word_count": len(" ".join(full_text).split()),
            },
            "structure": [{"title": f"Slide {s['slide']}", "content": s["text"]} for s in slides],
        }

    async def _parse_text(self, content: bytes, filename: str) -> Dict:
        """Parse plain text file."""
        text = content.decode("utf-8", errors="ignore")

        return {
            "text": text,
            "metadata": {
                "word_count": len(text.split()),
                "line_count": len(text.splitlines()),
                "char_count": len(text),
            },
            "structure": self._extract_structure([text]),
        }

    async def _parse_markdown(self, content: bytes, filename: str) -> Dict:
        """Parse Markdown file."""
        import markdown
        from bs4 import BeautifulSoup

        text = content.decode("utf-8", errors="ignore")
        html = markdown.markdown(text)
        soup = BeautifulSoup(html, "html.parser")
        plain_text = soup.get_text()

        # Extract headings
        headings = []
        for level in range(1, 7):
            for heading in soup.find_all(f"h{level}"):
                headings.append(
                    {
                        "text": heading.get_text(),
                        "level": level,
                    }
                )

        # Extract code blocks
        code_blocks = []
        for code in soup.find_all("code"):
            code_blocks.append(code.get_text())

        return {
            "text": plain_text,
            "markdown": text,
            "html": html,
            "headings": headings,
            "code_blocks": code_blocks,
            "metadata": {
                "word_count": len(plain_text.split()),
                "heading_count": len(headings),
                "code_block_count": len(code_blocks),
            },
            "structure": self._build_structure_from_headings(headings),
        }

    async def _parse_html(self, content: bytes, filename: str) -> Dict:
        """Parse HTML file."""
        from bs4 import BeautifulSoup

        html = content.decode("utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()

        text = soup.get_text(separator="\n", strip=True)

        # Extract headings
        headings = []
        for level in range(1, 7):
            for heading in soup.find_all(f"h{level}"):
                headings.append(
                    {
                        "text": heading.get_text(),
                        "level": level,
                    }
                )

        # Extract links
        links = []
        for link in soup.find_all("a", href=True):
            links.append(
                {
                    "text": link.get_text(),
                    "href": link["href"],
                }
            )

        return {
            "text": text,
            "html": html,
            "headings": headings,
            "links": links,
            "metadata": {
                "title": soup.title.string if soup.title else filename,
                "word_count": len(text.split()),
                "link_count": len(links),
            },
            "structure": self._build_structure_from_headings(headings),
        }

    async def _parse_image(self, content: bytes, filename: str) -> Dict:
        """Parse image using OCR."""
        if not self._ocr_enabled:
            return {
                "text": "",
                "metadata": {
                    "error": "OCR not available",
                    "filename": filename,
                },
                "structure": [],
            }

        import pytesseract
        from PIL import Image

        image = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(image)

        return {
            "text": text,
            "metadata": {
                "width": image.width,
                "height": image.height,
                "format": image.format,
                "word_count": len(text.split()),
            },
            "structure": self._extract_structure([text]),
        }

    def _extract_structure(self, text_parts: List[str]) -> List[Dict]:
        """Extract document structure from text."""
        structure = []
        full_text = "\n".join(text_parts)

        # Look for section patterns
        patterns = [
            (r"^(#{1,6})\s+(.+)$", "markdown_heading"),
            (r"^(\d+\.)\s+(.+)$", "numbered"),
            (r"^([A-Z][^.!?]*:)\s*$", "section_header"),
            (r"^(Chapter|Section|Part)\s+(\d+|[IVXLC]+)[\.:]\s*(.+)$", "chapter"),
        ]

        for line in full_text.splitlines():
            line = line.strip()
            if not line:
                continue

            for pattern, ptype in patterns:
                match = re.match(pattern, line, re.MULTILINE)
                if match:
                    structure.append(
                        {
                            "type": ptype,
                            "text": line,
                            "level": len(match.group(1)) if ptype == "markdown_heading" else 1,
                        }
                    )
                    break

        return structure

    def _build_structure_from_headings(self, headings: List[Dict]) -> List[Dict]:
        """Build hierarchical structure from headings."""
        if not headings:
            return []

        structure = []
        stack = [{"level": 0, "children": structure}]

        for heading in headings:
            level = heading.get("level", 1)
            node = {
                "title": heading["text"],
                "level": level,
                "children": [],
            }

            # Find parent
            while stack and stack[-1]["level"] >= level:
                stack.pop()

            if stack:
                stack[-1].get("children", structure).append(node)
            else:
                structure.append(node)

            stack.append(node)

        return structure

    def _get_heading_level(self, style_name: str) -> int:
        """Extract heading level from Word style name."""
        match = re.search(r"(\d+)", style_name)
        return int(match.group(1)) if match else 1

    async def chunk_text(
        self,
        text: str,
        chunk_size: int = 1000,
        overlap: int = 200,
        respect_sentences: bool = True,
    ) -> List[Dict]:
        """
        Split text into chunks for processing.

        Args:
            text: Input text
            chunk_size: Target chunk size in characters
            overlap: Overlap between chunks
            respect_sentences: Try to break at sentence boundaries

        Returns:
            List of chunks with metadata
        """
        if not text:
            return []

        chunks = []

        if respect_sentences:
            sentences = re.split(r"(?<=[.!?])\s+", text)
            current_chunk = []
            current_length = 0

            for sentence in sentences:
                sentence_length = len(sentence)

                if current_length + sentence_length > chunk_size and current_chunk:
                    chunk_text = " ".join(current_chunk)
                    chunks.append(
                        {
                            "text": chunk_text,
                            "start": len(" ".join(chunks)) if chunks else 0,
                            "length": len(chunk_text),
                        }
                    )

                    # Keep overlap
                    overlap_text = []
                    overlap_length = 0
                    for s in reversed(current_chunk):
                        if overlap_length + len(s) <= overlap:
                            overlap_text.insert(0, s)
                            overlap_length += len(s)
                        else:
                            break

                    current_chunk = overlap_text
                    current_length = overlap_length

                current_chunk.append(sentence)
                current_length += sentence_length

            if current_chunk:
                chunks.append(
                    {
                        "text": " ".join(current_chunk),
                        "start": len(" ".join([c["text"] for c in chunks])) if chunks else 0,
                        "length": len(" ".join(current_chunk)),
                    }
                )
        else:
            for i in range(0, len(text), chunk_size - overlap):
                chunk = text[i : i + chunk_size]
                chunks.append(
                    {
                        "text": chunk,
                        "start": i,
                        "length": len(chunk),
                    }
                )

        return chunks
